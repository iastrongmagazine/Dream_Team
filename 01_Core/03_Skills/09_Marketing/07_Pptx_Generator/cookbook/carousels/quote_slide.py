"""
Script to generate a square quote slide for presentations using python-pptx.
Includes branding colors, typography, and automatic text fitting.
"""

import sys
from pathlib import Path

# Tentative imports for pptx. These are often not available in all environments.
try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.util import Inches, Pt
except ImportError:
    # Print warning if pptx is missing but allow script to be linted/inspected
    Presentation = None
    RGBColor = None
    MSO_SHAPE = None
    PP_ALIGN = None
    MSO_ANCHOR = None
    Inches = None
    Pt = None
    print("Warning: python-pptx not found. Install with 'pip install python-pptx'")


def hex_to_rgb(hex_color: str):
    """Converts a hex color string (e.g., '#FFFFFF') to an RGBColor object."""
    h = hex_color.lstrip("#")
    if RGBColor is None:
        raise ImportError("RGBColor not available. Ensure python-pptx is installed.")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def main() -> None:
    """Generates a square quote slide with attribution and branding."""
    # === BRAND COLORS ===
    brand_bg = "1e1e2e"
    brand_text = "cdd6f4"
    brand_accent = "fab387"
    brand_accent_secondary = "89b4fa"
    brand_heading_font = "JetBrains Mono"
    brand_body_font = "Inter"

    # === CONTENT ===
    quote_text = "REPLACE"  # Max 120 chars, no quotation marks
    attribution = ""        # Optional, max 40 chars

    # === SETUP PRESENTATION ===
    if Presentation is None:
        print("Error: Presentation class not available. Ensure python-pptx is installed.")
        sys.exit(1)

    prs = Presentation()

    # Set slide size to 1080x1080 (Square for Instagram/FullHD)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(10)

    # Add blank slide
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # === BACKGROUND ===
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(brand_bg)

    # === DECORATIVE ELEMENTS ===
    # Accent bar at top
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.1)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = hex_to_rgb(brand_accent)
    top_bar.line.fill.background()

    # === QUOTE TEXT ===
    left = Inches(1)
    top = Inches(2.5)
    width = prs.slide_width - Inches(2)
    height = Inches(4)

    tx_box = slide.shapes.add_textbox(left, top, width, height)
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.text = f'"{quote_text}"'
    p.alignment = PP_ALIGN.CENTER
    p.font.name = brand_body_font
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(brand_text)

    # === ATTRIBUTION ===
    if attribution:
        attr_box = slide.shapes.add_textbox(
            Inches(1), Inches(6.5), prs.slide_width - Inches(2), Inches(1)
        )
        ap = attr_box.text_frame.paragraphs[0]
        ap.text = f"— {attribution}"
        ap.alignment = PP_ALIGN.CENTER
        ap.font.name = brand_heading_font
        ap.font.size = Pt(32)
        ap.font.color.rgb = hex_to_rgb(brand_accent_secondary)

    # === LOGO / BRANDING ===
    logo_box = slide.shapes.add_textbox(
        Inches(1), Inches(8.5), prs.slide_width - Inches(2), Inches(0.5)
    )
    lp = logo_box.text_frame.paragraphs[0]
    lp.text = "AI STRONG MAGAZINE"
    lp.alignment = PP_ALIGN.CENTER
    lp.font.name = brand_heading_font
    lp.font.size = Pt(18)
    lp.font.color.rgb = hex_to_rgb(brand_accent)

    # === SAVE ===
    output_path = Path("output_slide.pptx")
    prs.save(output_path)
    print(f"✅ Slide generated: {output_path.absolute()}")


if __name__ == "__main__":
    main()
