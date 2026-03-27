"""
83_Universal_Parser.py - Universal Doc Reader Elite
==================================================
Extrae información estructurada de archivos Office, PDF, PSD, e imágenes.
Producir salida Markdown premium optimizada para RAG/LLM.

Uso:
    python 83_Universal_Parser.py <file_path>
    python 83_Universal_Parser.py ./documentos/informe.pdf

Skills relacionadas:
    Skill 11_01: Universal Doc Reader Elite
    Skill 11_02: Batch Doc Processor Elite
"""

import os
import sys
import argparse
import io
from pathlib import Path
from datetime import datetime
from typing import Optional

# =============================================================================
# ARMOR LAYER - Configuración UTF-8
# =============================================================================
if sys.stdout.encoding != "utf-8":
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")

# =============================================================================
# IMPORTS - Con fallback para dependencias opcionales
# =============================================================================
try:
    import PyPDF2

    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from docx import Document

    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import pandas as pd

    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False

try:
    from pptx import Presentation

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from PIL import Image

    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False

try:
    from psd_tools import PSDImage

    PSD_AVAILABLE = True
except ImportError:
    PSD_AVAILABLE = False

# =============================================================================
# CONFIG PATHS
# =============================================================================
try:
    sys.path.insert(0, str(Path(__file__).parent))
    from config_paths import ROOT_DIR, KNOWLEDGE_DIR
except ImportError:
    ROOT_DIR = Path(__file__).resolve().parent.parent.parent
    KNOWLEDGE_DIR = ROOT_DIR / "03_Knowledge"

# =============================================================================
# BRANDING & COLORS
# =============================================================================
SCRIPT_NAME = "25_Universal_Parser.py"
SCRIPT_VERSION = "v2.0"
SKILL_ID = "31"
SKILL_NAME = "Universal Doc Reader Elite"


def print_branding():
    """Header premium del script."""
    print("─" * 60)
    print(f"  🦾 PersonalOS | Skill {SKILL_ID}: {SKILL_NAME}")
    print(f"  ⚡ {SCRIPT_NAME} {SCRIPT_VERSION}")
    print("─" * 60)
    print(f"  📅 {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    print(f"  📂 Root: {ROOT_DIR}")
    print("─" * 60)
    print()


def print_error(msg: str):
    print(f"  ❌ ERROR: {msg}", file=sys.stderr)


def print_success(msg: str):
    print(f"  ✅ {msg}")


def print_info(msg: str):
    print(f"  ℹ️  {msg}")


# =============================================================================
# PARSERS - Por formato
# =============================================================================


def parse_pdf(file_path: str) -> str:
    """Extrae texto de PDF."""
    if not PDF_AVAILABLE:
        return "ERROR: PyPDF2 no instalado. Ejecuta: pip install PyPDF2"

    text = []
    try:
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            metadata = reader.metadata
            text.append(f"## PDF: {Path(file_path).stem}\n")
            if metadata:
                text.append(f"**Autor:** {metadata.get('/Author', 'N/A')}")
                text.append(f"**Título:** {metadata.get('/Title', 'N/A')}")
                text.append(f"**Páginas:** {len(reader.pages)}\n")
                text.append("---\n\n")
            for i, page in enumerate(reader.pages):
                page_text = page.extract_text() or ""
                if page_text.strip():
                    text.append(f"### Página {i + 1}\n\n{page_text}\n\n")
        return "\n".join(text)
    except Exception as e:
        return f"ERROR parsing PDF: {e}"


def parse_docx(file_path: str) -> str:
    """Extrae texto y tablas de DOCX."""
    if not DOCX_AVAILABLE:
        return "ERROR: python-docx no instalado. Ejecuta: pip install python-docx"

    try:
        doc = Document(file_path)
        output = [f"## DOCX: {Path(file_path).stem}\n"]

        # Extraer párrafos
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        if paragraphs:
            output.append("\n".join(paragraphs))

        # Extraer tablas
        if doc.tables:
            output.append("\n\n### Tablas\n")
            for i, table in enumerate(doc.tables):
                output.append(f"#### Tabla {i + 1}\n")
                headers = [cell.text.strip() for cell in table.rows[0].cells]
                output.append("| " + " | ".join(headers) + " |")
                output.append("| " + " | ".join(["---"] * len(headers)) + " |")
                for row in table.rows[1:]:
                    cells = [cell.text.strip() for cell in row.cells]
                    output.append("| " + " | ".join(cells) + " |")
                output.append("")

        return "\n".join(output)
    except Exception as e:
        return f"ERROR parsing DOCX: {e}"


def parse_excel(file_path: str) -> str:
    """Extrae datos de XLSX, XLS y CSV."""
    if not PANDAS_AVAILABLE:
        return "ERROR: pandas no instalado. Ejecuta: pip install pandas openpyxl"

    try:
        output = [f"## Excel/CSV: {Path(file_path).stem}\n"]

        if file_path.lower().endswith(".csv"):
            df = pd.read_csv(file_path)
            output.append(f"**Filas:** {len(df)} | **Columnas:** {len(df.columns)}\n")
            output.append(df.to_markdown(index=False))
        else:
            dfs = pd.read_excel(file_path, sheet_name=None)
            for sheet_name, df in dfs.items():
                output.append(f"### Sheet: {sheet_name}")
                output.append(
                    f"**Filas:** {len(df)} | **Columnas:** {len(df.columns)}\n"
                )
                output.append(df.to_markdown(index=False))
                output.append("")

        return "\n".join(output)
    except Exception as e:
        return f"ERROR parsing Excel: {e}"


def parse_pptx(file_path: str) -> str:
    """Extrae contenido de PowerPoint."""
    if not PPTX_AVAILABLE:
        return "ERROR: python-pptx no instalado. Ejecuta: pip install python-pptx"

    try:
        prs = Presentation(file_path)
        output = [f"## PPTX: {Path(file_path).stem}\n"]
        output.append(f"**Total Slides:** {len(prs.slides)}\n\n")

        for i, slide in enumerate(prs.slides):
            output.append(f"### Slide {i + 1}\n")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    output.append(f"{shape.text}\n")
            output.append("---\n")
        return "\n".join(output)
    except Exception as e:
        return f"ERROR parsing PPTX: {e}"


def parse_psd(file_path: str) -> str:
    """Extrae metadata y estructura de capas de PSD."""
    if not PSD_AVAILABLE:
        return "ERROR: psd-tools no instalado. Ejecuta: pip install psd-tools"

    try:
        psd = PSDImage.open(file_path)
        output = [
            f"## PSD: {Path(file_path).stem}\n",
            f"**Dimensiones:** {psd.width} x {psd.height} px",
            f"**Modo:** {psd.color_mode}",
            "",
            "### Capas",
            "",
        ]
        layer_count = 0
        for layer in psd:
            layer_count += 1
            output.append(f"- **{layer.name}** ({layer.kind})")
        output.append(f"\n**Total capas:** {layer_count}")
        return "\n".join(output)
    except Exception as e:
        return f"ERROR parsing PSD: {e}"


def parse_image_meta(file_path: str) -> str:
    """Extrae metadata de imágenes."""
    if not PIL_AVAILABLE:
        return "ERROR: Pillow no instalado. Ejecuta: pip install Pillow"

    try:
        with Image.open(file_path) as img:
            output = [
                f"## Imagen: {Path(file_path).stem}\n",
                f"**Formato:** {img.format}",
                f"**Modo:** {img.mode}",
                f"**Dimensiones:** {img.width} x {img.height} px",
            ]
            if hasattr(img, "_getexif") and img._getexif():
                output.append("\n### EXIF Data")
                for tag_id, value in img._getexif().items():
                    output.append(f"- {tag_id}: {value}")
            return "\n".join(output)
    except Exception as e:
        return f"ERROR parsing imagen: {e}"


def parse_markdown(file_path: str) -> str:
    """Lee archivos Markdown."""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    except Exception as e:
        return f"ERROR leyendo MD: {e}"


# =============================================================================
# MAPPER - Extensión → Parser
# =============================================================================
PARSER_MAP = {
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".doc": parse_docx,
    ".xlsx": parse_excel,
    ".xls": parse_excel,
    ".csv": parse_excel,
    ".pptx": parse_pptx,
    ".psd": parse_psd,
    ".psb": parse_psd,
    ".tiff": parse_image_meta,
    ".tif": parse_image_meta,
    ".jpg": parse_image_meta,
    ".jpeg": parse_image_meta,
    ".png": parse_image_meta,
    ".gif": parse_image_meta,
    ".webp": parse_image_meta,
    ".md": parse_markdown,
    ".txt": parse_markdown,
}


def main():
    parser = argparse.ArgumentParser(
        description=f"PersonalOS Skill {SKILL_ID} - {SKILL_NAME}",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Ejemplos:
  python 83_Universal_Parser.py ./docs/informe.pdf
  python 83_Universal_Parser.py ./datos/ventas.xlsx
  python 83_Universal_Parser.py "./presentación cliente.pptx"
        """,
    )
    parser.add_argument("file", help="Ruta al archivo a procesar")
    parser.add_argument(
        "--output", "-o", help="Guardar resultado en archivo (opcional)", default=None
    )
    args = parser.parse_args()

    # Validar archivo
    file_path = Path(args.file).resolve()
    if not file_path.exists():
        print_error(f"Archivo no encontrado: {file_path}")
        sys.exit(1)

    ext = file_path.suffix.lower()

    # Branding
    print_branding()
    print_info(f"Procesando: {file_path.name} ({ext})")

    # Obtener parser
    parse_func = PARSER_MAP.get(ext)
    if not parse_func:
        print_error(f"Formato no soportado: {ext}")
        print_info(f"Soportados: {', '.join(PARSER_MAP.keys())}")
        sys.exit(1)

    # Extraer contenido
    content = parse_func(str(file_path))

    # Output
    print()
    if args.output:
        output_path = Path(args.output)
        output_path.write_text(content, encoding="utf-8")
        print_success(f"Guardado en: {output_path}")
    else:
        print(content)

    print()
    print("─" * 60)
    print_success(f"Extracción completa: {file_path.name}")
    print("─" * 60)


if __name__ == "__main__":
    main()
